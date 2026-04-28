"""Unit-test suite for `pptx.chart.quick_layouts`."""

from __future__ import annotations

import pytest

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.chart.quick_layouts import (
    QUICK_LAYOUTS,
    apply_quick_layout,
    layout_names,
)
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches


def _make_column_chart(series=("S1", "S2")):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    data = CategoryChartData()
    data.categories = ["A", "B", "C"]
    for name in series:
        data.add_series(name, (1, 2, 3))
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(4),
        data,
    )
    return gframe.chart


def _make_pie_chart():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    data = CategoryChartData()
    data.categories = ["A", "B"]
    data.add_series("S1", (1, 2))
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(1), Inches(1), Inches(6), Inches(4), data
    )
    return gframe.chart


class DescribeApplyQuickLayout:
    def it_applies_a_named_layout_with_legend_on_the_right(self):
        chart = _make_column_chart()

        apply_quick_layout(chart, "title_legend_right")

        assert chart.has_title is True
        assert chart.has_legend is True
        assert chart.legend.position == XL_LEGEND_POSITION.RIGHT
        assert chart.value_axis.has_major_gridlines is True
        assert chart.value_axis.has_minor_gridlines is False

    def it_applies_minimal_layout(self):
        chart = _make_column_chart()

        apply_quick_layout(chart, "minimal")

        assert chart.has_title is False
        assert chart.has_legend is False
        assert chart.value_axis.has_major_gridlines is False
        assert chart.value_axis.has_minor_gridlines is False

    def it_applies_dense_layout_with_minor_gridlines(self):
        chart = _make_column_chart()

        apply_quick_layout(chart, "dense")

        assert chart.has_title is True
        assert chart.has_legend is True
        assert chart.value_axis.has_major_gridlines is True
        assert chart.value_axis.has_minor_gridlines is True

    def it_can_set_title_text_via_spec(self):
        chart = _make_column_chart()

        apply_quick_layout(chart, {"title_text": "Q4 Revenue"})

        assert chart.has_title is True
        assert chart.chart_title.text_frame.text == "Q4 Revenue"

    def it_can_set_axis_titles_via_spec(self):
        chart = _make_column_chart()

        apply_quick_layout(
            chart,
            {
                "category_axis_title_text": "Quarter",
                "value_axis_title_text": "Revenue (USD)",
            },
        )

        assert chart.category_axis.has_title is True
        assert chart.category_axis.axis_title.text_frame.text == "Quarter"
        assert chart.value_axis.has_title is True
        assert chart.value_axis.axis_title.text_frame.text == "Revenue (USD)"

    def it_supports_composition(self):
        chart = _make_column_chart()

        apply_quick_layout(chart, "title_legend_right")
        apply_quick_layout(chart, {"has_minor_gridlines": True})

        assert chart.has_legend is True  # preserved from first layout
        assert chart.legend.position == XL_LEGEND_POSITION.RIGHT
        assert chart.value_axis.has_minor_gridlines is True

    def it_silently_skips_axis_keys_on_pie_charts(self):
        chart = _make_pie_chart()

        # Pie has no value/category axis; the call should still succeed.
        apply_quick_layout(chart, "title_axes_legend_right")

        assert chart.has_title is True
        assert chart.has_legend is True

    def it_does_not_touch_legend_position_when_legend_is_off(self):
        chart = _make_column_chart()
        chart.has_legend = False

        # `title_no_legend` doesn't set legend_position; legend should stay off.
        apply_quick_layout(chart, "title_no_legend")

        assert chart.has_legend is False
        assert chart.legend is None

    def it_raises_for_unknown_layout_name(self):
        chart = _make_column_chart()

        with pytest.raises(ValueError, match="unknown quick layout"):
            apply_quick_layout(chart, "not_a_real_layout")

    def it_raises_for_invalid_layout_type(self):
        chart = _make_column_chart()

        with pytest.raises(TypeError, match="must be a name or spec"):
            apply_quick_layout(chart, 42)

    def it_is_exposed_as_a_method_on_Chart(self):
        chart = _make_column_chart()

        chart.apply_quick_layout("title_legend_bottom")

        assert chart.legend.position == XL_LEGEND_POSITION.BOTTOM


class DescribeLayoutNames:
    def it_returns_all_built_in_names(self):
        names = layout_names()
        assert isinstance(names, tuple)
        assert set(names) == set(QUICK_LAYOUTS.keys())

    def it_preserves_declaration_order(self):
        # Make sure the declaration-order contract is testable; first
        # entry should be the most generic preset.
        assert layout_names()[0] == "title_legend_right"
