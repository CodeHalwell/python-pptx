"""Baseline round-trip regression suite for the Phase 1 hygiene fixes.

Every later roadmap phase is expected to add a parametrized test here that
covers the new feature it ships, so that PowerPoint-authored decks using that
feature aren't silently corrupted by `Presentation(...).save()`.
"""

from __future__ import annotations

import pytest

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from .round_trip import assert_round_trip, round_trip_diff


class DescribeRoundTrip:
    def it_round_trips_an_empty_deck(self):
        assert_round_trip(Presentation())

    def it_round_trips_a_deck_with_a_blank_slide(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        assert_round_trip(prs)

    def it_round_trips_a_deck_with_text(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text_frame.text = "Round-trip"
        assert_round_trip(prs)

    def it_round_trips_a_deck_with_an_autoshape(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))
        assert_round_trip(prs)

    def it_round_trips_a_deck_with_an_explicit_font_color(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tf = slide.shapes.title.text_frame
        tf.text = "Colored"
        tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xC0, 0x10, 0x40)
        assert_round_trip(prs)

    def it_round_trips_a_deck_with_an_explicit_line_color(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))
        shape.line.color.rgb = RGBColor(0x10, 0xA0, 0x10)
        assert_round_trip(prs)


class DescribeNonMutatingColorReadsRoundTrip:
    """Reading `Font.color` properties must not corrupt the round-trip."""

    @pytest.mark.parametrize(
        "prop", ["type", "rgb", "theme_color", "brightness"]
    )
    def it_does_not_alter_a_decks_xml_after_a_color_read(self, prop):
        # -- compose a fresh deck and snapshot the post-save bytes --
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tf = slide.shapes.title.text_frame
        tf.text = "Inheritance"
        font = tf.paragraphs[0].runs[0].font

        # -- just reading should not mutate the underlying XML or perturb the
        # -- save → open → save invariant --
        getattr(font.color, prop)
        getattr(font.color, prop)

        assert round_trip_diff(prs) == {}

    @pytest.mark.parametrize(
        "prop", ["type", "rgb", "theme_color", "brightness"]
    )
    def it_does_not_alter_a_shapes_xml_after_a_line_color_read(self, prop):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))

        getattr(shape.line.color, prop)
        getattr(shape.line.color, prop)

        assert round_trip_diff(prs) == {}


class DescribeShapeIdAllocationRoundTrip:
    """The cached-cursor allocator must not collide with existing ids."""

    def it_does_not_emit_duplicate_shape_ids(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for _ in range(20):
            slide.shapes.add_shape(1, Inches(0), Inches(0), Pt(10), Pt(10))

        ids = [int(v) for v in slide.shapes._spTree.xpath("//@id") if v.isdigit()]
        assert len(ids) == len(set(ids)), "duplicate shape ids: %r" % ids
        assert_round_trip(prs)


class DescribeAnimationsRoundTrip:
    """Phase 5 animation extensions must survive a save→open→save cycle."""

    def it_round_trips_a_motion_path_animation(self):
        from pptx.animation import MotionPath

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))
        MotionPath.line(slide, shape, Inches(2), Inches(0))
        MotionPath.custom(slide, shape, "M 0 0 L 0.5 0.25 E")
        assert_round_trip(prs)

    def it_round_trips_a_sequenced_animation_chain(self):
        from pptx.animation import Emphasis, Entrance

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        a = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))
        b = slide.shapes.add_shape(1, Inches(1), Inches(2.5), Inches(2), Inches(1))
        c = slide.shapes.add_shape(1, Inches(1), Inches(4), Inches(2), Inches(1))
        with slide.animations.sequence():
            Entrance.fade(slide, a)
            Entrance.fly_in(slide, b)
            Emphasis.pulse(slide, c)
        assert_round_trip(prs)

    def it_round_trips_a_by_paragraph_entrance(self):
        from pptx.animation import Entrance

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        tf = tb.text_frame
        tf.text = "alpha"
        tf.add_paragraph().text = "beta"
        tf.add_paragraph().text = "gamma"
        Entrance.fade(slide, tf, by_paragraph=True)
        assert_round_trip(prs)


class DescribeThemeWriterRoundTrip:
    """Phase 7 writable theme: changes must persist across save/open/save."""

    def it_round_trips_an_overridden_accent_color(self):
        from pptx.enum.dml import MSO_THEME_COLOR

        prs = Presentation()
        prs.theme.colors[MSO_THEME_COLOR.ACCENT_1] = RGBColor(0xFF, 0x66, 0x00)
        assert_round_trip(prs)

    def it_round_trips_overridden_fonts(self):
        prs = Presentation()
        prs.theme.fonts.major = "Inter"
        prs.theme.fonts.minor = "Inter"
        assert_round_trip(prs)
