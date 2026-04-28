"""Unit-test suite for `pptx.render`.

LibreOffice is mocked out: these tests cover argument plumbing, error
handling, and selection of slide indexes — not the actual rendering.
"""

from __future__ import annotations

import tempfile
from pathlib import Path
from unittest.mock import patch

import pytest

from pptx import Presentation
from pptx.render import (
    DEFAULT_TIMEOUT_SECONDS,
    ThumbnailRendererError,
    ThumbnailRendererUnavailable,
    render_slide_thumbnail,
    render_slide_thumbnails,
)


@pytest.fixture
def two_slide_prs():
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    prs.slides.add_slide(prs.slide_layouts[5])
    return prs


def _fake_soffice_run(work_dir: Path, *, num_slides=2, exit_code=0, stderr=b""):
    """Drop fake `*.png` files in the work_dir and return a stub CompletedProcess."""
    from subprocess import CompletedProcess

    def _runner(soffice_bin, deck_path, out_dir, timeout):
        if exit_code == 0:
            for i in range(num_slides):
                (Path(out_dir) / f"slide{i}.png").write_bytes(b"\x89PNG\r\n\x1a\n%d" % i)
        return CompletedProcess(args=[], returncode=exit_code, stdout=b"", stderr=stderr)

    return _runner


def _fake_soffice_run_with_names(work_dir: Path, names):
    """Variant fake runner that writes specific filenames in a given order."""
    from subprocess import CompletedProcess

    def _runner(soffice_bin, deck_path, out_dir, timeout):
        for n in names:
            (Path(out_dir) / n).write_bytes(b"\x89PNG\r\n\x1a\n")
        return CompletedProcess(args=[], returncode=0, stdout=b"", stderr=b"")

    return _runner


class DescribeRenderSlideThumbnails:
    def it_returns_paths_for_every_slide_by_default(self, tmp_path, two_slide_prs):
        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=_fake_soffice_run(tmp_path)
        ):
            paths = render_slide_thumbnails(two_slide_prs, out_dir=tmp_path)

        assert len(paths) == 2
        assert all(p.suffix == ".png" for p in paths)

    def it_returns_only_requested_indexes(self, tmp_path, two_slide_prs):
        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=_fake_soffice_run(tmp_path)
        ):
            paths = render_slide_thumbnails(
                two_slide_prs, out_dir=tmp_path, slide_indexes=[1]
            )

        assert len(paths) == 1
        assert paths[0].name.endswith("1.png")

    def it_raises_when_index_is_out_of_range(self, tmp_path, two_slide_prs):
        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=_fake_soffice_run(tmp_path)
        ):
            with pytest.raises(IndexError, match="out of range"):
                render_slide_thumbnails(
                    two_slide_prs, out_dir=tmp_path, slide_indexes=[5]
                )

    def it_returns_bytes_when_asked(self, tmp_path, two_slide_prs):
        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=_fake_soffice_run(tmp_path)
        ):
            data = render_slide_thumbnails(
                two_slide_prs, out_dir=tmp_path, return_bytes=True
            )

        assert all(isinstance(d, bytes) for d in data)
        assert data[0].startswith(b"\x89PNG")

    def it_raises_unavailable_when_soffice_is_missing(self, two_slide_prs):
        with patch("pptx.render.shutil.which", return_value=None):
            with pytest.raises(ThumbnailRendererUnavailable, match="install LibreOffice"):
                render_slide_thumbnails(two_slide_prs)

    def it_raises_when_soffice_exits_non_zero(self, tmp_path, two_slide_prs):
        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice",
            side_effect=_fake_soffice_run(tmp_path, exit_code=1, stderr=b"boom"),
        ):
            with pytest.raises(ThumbnailRendererError, match="status 1"):
                render_slide_thumbnails(two_slide_prs, out_dir=tmp_path)

    def it_raises_when_soffice_produces_no_pngs(self, tmp_path, two_slide_prs):
        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=_fake_soffice_run(tmp_path, num_slides=0)
        ):
            with pytest.raises(ThumbnailRendererError, match="no PNG output"):
                render_slide_thumbnails(two_slide_prs, out_dir=tmp_path)

    def it_passes_the_default_timeout_through(self, tmp_path, two_slide_prs):
        captured = {}

        def _capture(soffice_bin, deck_path, out_dir, timeout):
            captured["timeout"] = timeout
            return _fake_soffice_run(tmp_path)(soffice_bin, deck_path, out_dir, timeout)

        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=_capture
        ):
            render_slide_thumbnails(two_slide_prs, out_dir=tmp_path)

        assert captured["timeout"] == DEFAULT_TIMEOUT_SECONDS

    def it_honours_custom_soffice_bin(self, tmp_path, two_slide_prs):
        captured = {}

        def _capture(soffice_bin, deck_path, out_dir, timeout):
            captured["bin"] = soffice_bin
            return _fake_soffice_run(tmp_path)(soffice_bin, deck_path, out_dir, timeout)

        with patch("pptx.render.shutil.which", return_value="/opt/libre/soffice"), patch(
            "pptx.render._run_soffice", side_effect=_capture
        ):
            render_slide_thumbnails(
                two_slide_prs, out_dir=tmp_path, soffice_bin="/opt/libre/soffice"
            )

        assert captured["bin"] == "/opt/libre/soffice"


class DescribeRenderSlideThumbnail:
    def it_renders_the_specific_slide(self, tmp_path, two_slide_prs):
        slide = two_slide_prs.slides[1]

        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=_fake_soffice_run(tmp_path)
        ):
            target = tmp_path / "out" / "slide.png"
            path = render_slide_thumbnail(slide, out_path=target)

        assert path == target
        assert path.exists()

    def it_returns_bytes_when_asked(self, tmp_path, two_slide_prs):
        slide = two_slide_prs.slides[0]

        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=_fake_soffice_run(tmp_path)
        ):
            data = render_slide_thumbnail(slide, return_bytes=True)

        assert isinstance(data, bytes)
        assert data.startswith(b"\x89PNG")


class DescribePresentationConvenienceMethods:
    def it_exposes_render_thumbnails_on_presentation(self, tmp_path, two_slide_prs):
        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=_fake_soffice_run(tmp_path)
        ):
            paths = two_slide_prs.render_thumbnails(out_dir=tmp_path)

        assert len(paths) == 2

    def it_exposes_render_thumbnail_on_slide(self, tmp_path, two_slide_prs):
        slide = two_slide_prs.slides[0]

        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=_fake_soffice_run(tmp_path)
        ):
            path = slide.render_thumbnail(out_path=tmp_path / "s.png")

        assert path.exists()


class DescribeNaturalSort:
    """Regression tests for the lexicographic-sort bug on 10+ slide decks."""

    def it_orders_double_digit_slide_numbers_correctly(self, tmp_path):
        from pptx import Presentation
        from pptx.render import _natural_sort_key

        prs = Presentation()
        for _ in range(11):
            prs.slides.add_slide(prs.slide_layouts[5])

        # Names mirror what LibreOffice emits: deck-1.png .. deck-11.png
        names = [f"deck-{i}.png" for i in range(1, 12)]
        runner = _fake_soffice_run_with_names(tmp_path, names)

        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=runner
        ):
            paths = render_slide_thumbnails(prs, out_dir=tmp_path)

        # With natural sort, slide #2 (index 1) is the second entry — not #10.
        assert [p.name for p in paths] == names

        # Spot-check the key for completeness.
        assert _natural_sort_key(Path("deck-2.png")) < _natural_sort_key(
            Path("deck-10.png")
        )

    def it_returns_the_correct_slide_for_high_indexes(self, tmp_path):
        """`slide_indexes=[9]` must hit `deck-10.png`, not `deck-2.png`."""
        from pptx import Presentation

        prs = Presentation()
        for _ in range(11):
            prs.slides.add_slide(prs.slide_layouts[5])

        names = [f"deck-{i}.png" for i in range(1, 12)]
        runner = _fake_soffice_run_with_names(tmp_path, names)

        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=runner
        ):
            paths = render_slide_thumbnails(
                prs, out_dir=tmp_path, slide_indexes=[9]
            )

        assert paths[0].name == "deck-10.png"


class DescribePngFiltering:
    """`render_slide_thumbnails` must ignore PNGs already in `out_dir`."""

    def it_ignores_preexisting_pngs_in_the_output_directory(self, tmp_path, two_slide_prs):
        # Drop a stray PNG in the directory before rendering — perhaps
        # left over from a previous job or simply present in a shared
        # artifacts folder.
        (tmp_path / "stale.png").write_bytes(b"\x89PNG\r\n\x1a\nstale")

        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=_fake_soffice_run(tmp_path)
        ):
            paths = render_slide_thumbnails(two_slide_prs, out_dir=tmp_path)

        names = {p.name for p in paths}
        assert "stale.png" not in names
        assert len(paths) == 2  # two slides, not three

    def it_keeps_indexes_aligned_when_strays_are_present(self, tmp_path, two_slide_prs):
        # A stray that would lex-sort *between* the two genuine outputs.
        (tmp_path / "slide05.png").write_bytes(b"\x89PNG\r\n\x1a\nstale")

        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=_fake_soffice_run(tmp_path)
        ):
            paths = render_slide_thumbnails(
                two_slide_prs, out_dir=tmp_path, slide_indexes=[1]
            )

        # Without the filter we'd risk returning slide05.png; with it,
        # index 1 maps to the genuine second render (slide1.png from the
        # fake runner, which writes slide{0,1}.png).
        assert paths[0].name == "slide1.png"


class DescribeRenderSlideThumbnailCleanup:
    """The single-slide renderer must not leak temp directories."""

    def it_cleans_the_temp_dir_when_out_path_is_given(self, tmp_path, two_slide_prs):
        slide = two_slide_prs.slides[0]
        created_dirs = []

        real_tmp = tempfile.TemporaryDirectory

        def _spy_tmp(*args, **kwargs):
            ctx = real_tmp(*args, **kwargs)
            created_dirs.append(ctx.name)
            return ctx

        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=_fake_soffice_run(tmp_path)
        ), patch("pptx.render.tempfile.TemporaryDirectory", side_effect=_spy_tmp):
            path = render_slide_thumbnail(slide, out_path=tmp_path / "s.png")

        assert path.exists()
        # Every temp dir we created in the spy must have been cleaned up.
        assert created_dirs, "expected at least one TemporaryDirectory"
        for d in created_dirs:
            assert not Path(d).exists(), f"temp dir {d} not cleaned"

    def it_cleans_the_temp_dir_when_no_out_path_is_given(self, tmp_path, two_slide_prs):
        slide = two_slide_prs.slides[0]
        created_dirs = []

        real_tmp = tempfile.TemporaryDirectory

        def _spy_tmp(*args, **kwargs):
            ctx = real_tmp(*args, **kwargs)
            created_dirs.append(ctx.name)
            return ctx

        with patch("pptx.render.shutil.which", return_value="/usr/bin/soffice"), patch(
            "pptx.render._run_soffice", side_effect=_fake_soffice_run(tmp_path)
        ), patch("pptx.render.tempfile.TemporaryDirectory", side_effect=_spy_tmp):
            path = render_slide_thumbnail(slide)

        try:
            assert path.exists()
            assert path.suffix == ".png"
            for d in created_dirs:
                assert not Path(d).exists(), f"temp dir {d} not cleaned"
        finally:
            # Caller owns this NamedTemporaryFile; clean up after the test.
            path.unlink(missing_ok=True)
