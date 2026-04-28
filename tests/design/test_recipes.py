"""Unit-test suite for :mod:`pptx.design.recipes`."""

from __future__ import annotations

import os

import pytest

from pptx import Presentation
from pptx.design.recipes import (
    bullet_slide,
    image_hero_slide,
    kpi_slide,
    quote_slide,
    title_slide,
)
from pptx.design.tokens import DesignTokens
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


_TEST_IMAGE = os.path.join(
    os.path.dirname(__file__), "..", "test_files", "monty-truth.png"
)


@pytest.fixture
def prs():
    return Presentation()


@pytest.fixture
def tokens():
    return DesignTokens.from_dict(
        {
            "palette": {
                "primary": "#3C2F80",
                "neutral": "#222222",
                "muted": "#777777",
                "surface": "#F4F2FB",
                "on_primary": "#FFFFFF",
                "positive": "#00853E",
                "negative": "#B00020",
            },
            "typography": {
                "heading": {"family": "Inter", "size": Pt(36), "bold": True},
                "body": {"family": "Inter", "size": Pt(16)},
            },
            "shadows": {
                "card": {
                    "blur_radius": Pt(8),
                    "distance": Pt(2),
                    "direction": 90.0,
                    "color": RGBColor(0, 0, 0),
                    "alpha": 0.2,
                },
            },
        }
    )


def _runs(slide):
    out = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                out.append(r)
    return out


def _paragraph_texts(slide):
    return [
        p.text
        for sh in slide.shapes
        if sh.has_text_frame
        for p in sh.text_frame.paragraphs
    ]


class DescribeTitleSlide:
    def it_appends_a_slide_to_the_presentation(self, prs):
        before = len(prs.slides)
        slide = title_slide(prs, title="Hello")
        assert len(prs.slides) == before + 1
        assert slide is prs.slides[-1]

    def it_writes_the_title_text(self, prs):
        slide = title_slide(prs, title="Q4 Review")
        assert any("Q4 Review" in t for t in _paragraph_texts(slide))

    def it_writes_the_subtitle_when_provided(self, prs):
        slide = title_slide(prs, title="T", subtitle="April 2026")
        assert any("April 2026" in t for t in _paragraph_texts(slide))

    def it_omits_the_subtitle_when_not_provided(self, prs):
        slide = title_slide(prs, title="T")
        # exactly one textbox: the title
        textboxes = [s for s in slide.shapes if s.has_text_frame]
        assert len(textboxes) == 1

    def it_applies_the_token_palette_to_the_title(self, prs, tokens):
        slide = title_slide(prs, title="Hi", tokens=tokens)
        runs = _runs(slide)
        assert runs[0].font.color.rgb == RGBColor(0x3C, 0x2F, 0x80)
        assert runs[0].font.name == "Inter"

    def it_centers_the_title(self, prs):
        slide = title_slide(prs, title="Hi")
        title_para = next(
            p for s in slide.shapes if s.has_text_frame for p in s.text_frame.paragraphs
        )
        assert title_para.alignment == PP_ALIGN.CENTER

    def it_applies_an_optional_transition(self, prs):
        from pptx.enum.presentation import MSO_TRANSITION_TYPE

        slide = title_slide(prs, title="Hi", transition="morph")
        assert slide.transition.kind == MSO_TRANSITION_TYPE.MORPH

    def it_raises_on_unknown_transition(self, prs):
        with pytest.raises(ValueError):
            title_slide(prs, title="Hi", transition="not_a_transition")


class DescribeBulletSlide:
    def it_writes_one_paragraph_per_bullet_plus_title(self, prs):
        slide = bullet_slide(prs, title="T", bullets=["a", "b", "c"])
        body_box = [s for s in slide.shapes if s.has_text_frame][1]
        # 3 bullets => 3 paragraphs
        assert len(body_box.text_frame.paragraphs) == 3

    def it_prefixes_a_bullet_glyph(self, prs):
        slide = bullet_slide(prs, title="T", bullets=["alpha"])
        assert any("•" in t and "alpha" in t for t in _paragraph_texts(slide))

    def it_applies_token_typography_to_bullets(self, prs, tokens):
        slide = bullet_slide(prs, title="T", bullets=["x"], tokens=tokens)
        runs = _runs(slide)
        # last run is the bullet
        assert runs[-1].font.name == "Inter"
        assert runs[-1].font.size == Pt(16)


class DescribeKpiSlide:
    def it_renders_one_card_per_kpi(self, prs):
        slide = kpi_slide(
            prs,
            title="Metrics",
            kpis=[
                {"label": "ARR", "value": "$182M", "delta": +0.27},
                {"label": "NDR", "value": "131%", "delta": -0.03},
            ],
        )
        cards = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        assert len(cards) == 2

    def it_writes_label_value_and_delta_text(self, prs):
        slide = kpi_slide(
            prs,
            title="Metrics",
            kpis=[{"label": "ARR", "value": "$182M", "delta": +0.27}],
        )
        texts = " ".join(_paragraph_texts(slide))
        assert "ARR" in texts
        assert "$182M" in texts
        assert "27%" in texts

    def it_handles_zero_kpis(self, prs):
        slide = kpi_slide(prs, title="None today", kpis=[])
        # only the title remains
        textboxes = [s for s in slide.shapes if s.has_text_frame]
        assert len(textboxes) == 1

    def it_tints_positive_and_negative_deltas_differently(self, prs, tokens):
        slide = kpi_slide(
            prs,
            title="Metrics",
            kpis=[
                {"label": "Up", "value": "1", "delta": +0.1},
                {"label": "Down", "value": "1", "delta": -0.1},
            ],
            tokens=tokens,
        )
        delta_runs = [
            run
            for sh in slide.shapes
            if sh.has_text_frame
            for p in sh.text_frame.paragraphs
            for run in p.runs
            if "%" in run.text
        ]
        colors = {run.font.color.rgb for run in delta_runs}
        assert RGBColor(0x00, 0x85, 0x3E) in colors
        assert RGBColor(0xB0, 0x00, 0x20) in colors

    def it_applies_the_card_shadow_token(self, prs, tokens):
        slide = kpi_slide(
            prs,
            title="Metrics",
            kpis=[{"label": "ARR", "value": "$182M"}],
            tokens=tokens,
        )
        card = next(
            s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        )
        assert card.shadow.blur_radius == Pt(8)
        assert card.shadow.distance == Pt(2)


class DescribeQuoteSlide:
    def it_wraps_the_quote_in_curly_quotes(self, prs):
        slide = quote_slide(prs, quote="Hello world")
        texts = _paragraph_texts(slide)
        assert any("“Hello world”" in t for t in texts)

    def it_renders_the_attribution_when_provided(self, prs):
        slide = quote_slide(prs, quote="Q", attribution="Ada")
        texts = _paragraph_texts(slide)
        assert any("Ada" in t for t in texts)

    def it_omits_the_attribution_when_missing(self, prs):
        slide = quote_slide(prs, quote="Q")
        textboxes = [s for s in slide.shapes if s.has_text_frame]
        assert len(textboxes) == 1


class DescribeImageHeroSlide:
    def it_adds_a_picture_at_full_slide_extent(self, prs):
        slide = image_hero_slide(prs, title="Demo", image=_TEST_IMAGE)
        pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pics) == 1
        assert pics[0].left == 0
        assert pics[0].top == 0
        assert pics[0].width == prs.slide_width
        assert pics[0].height == prs.slide_height

    def it_renders_the_caption_when_provided(self, prs):
        slide = image_hero_slide(
            prs, title="T", image=_TEST_IMAGE, caption="A short caption"
        )
        assert any("A short caption" in t for t in _paragraph_texts(slide))

    def it_uses_the_token_primary_color_for_the_band(self, prs, tokens):
        slide = image_hero_slide(
            prs, title="T", image=_TEST_IMAGE, tokens=tokens
        )
        bands = [
            s
            for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        ]
        assert bands
        assert bands[0].fill.fore_color.rgb == RGBColor(0x3C, 0x2F, 0x80)
