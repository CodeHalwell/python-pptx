"""Unit-test suite for the high-level `pptx.theme` module."""

from __future__ import annotations

import pytest

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.theme import Theme, ThemeColors, ThemeFonts


@pytest.fixture
def prs():
    return Presentation()


@pytest.fixture
def theme(prs) -> Theme:
    return prs.theme


class DescribeThemeColors:
    def it_reads_a_theme_color_slot(self, theme):
        rgb = theme.colors[MSO_THEME_COLOR.ACCENT_1]
        assert isinstance(rgb, RGBColor)

    def it_writes_a_theme_color_slot(self, theme):
        theme.colors[MSO_THEME_COLOR.ACCENT_2] = RGBColor(0x12, 0x34, 0x56)
        assert theme.colors[MSO_THEME_COLOR.ACCENT_2] == RGBColor(0x12, 0x34, 0x56)

    def it_writes_to_the_canonical_slot_for_an_alias(self, theme):
        # bg1 is an alias for lt1 — the write should land on lt1
        theme.colors[MSO_THEME_COLOR.BACKGROUND_1] = RGBColor(0xFF, 0x00, 0xFF)
        assert theme.colors[MSO_THEME_COLOR.BACKGROUND_1] == RGBColor(0xFF, 0x00, 0xFF)

    def it_raises_when_assigning_a_non_rgb_value(self, theme):
        with pytest.raises(TypeError):
            theme.colors[MSO_THEME_COLOR.ACCENT_1] = "purple"  # type: ignore[assignment]

    def it_raises_when_assigning_with_a_non_theme_color_key(self, theme):
        with pytest.raises(TypeError):
            theme.colors["accent1"] = RGBColor(0, 0, 0)  # type: ignore[index]

    def it_inserts_a_missing_slot_at_the_schema_defined_position(self, prs):
        from pptx.oxml.ns import qn

        clr_scheme = prs.theme._theme_elm.find(  # type: ignore[attr-defined]
            f"{qn('a:themeElements')}/{qn('a:clrScheme')}"
        )
        # Remove accent3 to force the writer's missing-slot path
        accent3 = clr_scheme.find(qn("a:accent3"))
        clr_scheme.remove(accent3)

        prs.theme.colors[MSO_THEME_COLOR.ACCENT_3] = RGBColor(0x12, 0x34, 0x56)

        order = [c.tag.rsplit("}", 1)[-1] for c in clr_scheme]
        # Schema-required ordering must be preserved
        assert order.index("accent2") + 1 == order.index("accent3")
        assert order.index("accent3") + 1 == order.index("accent4")

    def it_overwrites_a_non_srgb_color_child(self, prs):
        # Replace the existing srgbClr in the accent3 slot with sysClr,
        # then verify the writer drops the sysClr and replaces with srgbClr.
        from lxml import etree
        from pptx.oxml.ns import qn

        slot = prs.theme._theme_elm.find(  # type: ignore[attr-defined]
            f"{qn('a:themeElements')}/{qn('a:clrScheme')}/{qn('a:accent3')}"
        )
        for child in list(slot):
            slot.remove(child)
        sys_clr = etree.SubElement(slot, qn("a:sysClr"))
        sys_clr.set("val", "windowText")
        sys_clr.set("lastClr", "000000")

        prs.theme.colors[MSO_THEME_COLOR.ACCENT_3] = RGBColor(0xAA, 0xBB, 0xCC)
        assert slot.find(qn("a:sysClr")) is None
        assert slot.find(qn("a:srgbClr")).get("val") == "AABBCC"


class DescribeThemeFonts:
    def it_reads_major_and_minor(self, theme):
        assert isinstance(theme.fonts.major, str)
        assert isinstance(theme.fonts.minor, str)

    def it_writes_major_typeface(self, theme):
        theme.fonts.major = "Inter"
        assert theme.fonts.major == "Inter"

    def it_writes_minor_typeface(self, theme):
        theme.fonts.minor = "Source Sans Pro"
        assert theme.fonts.minor == "Source Sans Pro"

    def it_rejects_empty_typeface(self, theme):
        with pytest.raises(TypeError):
            theme.fonts.major = ""


class DescribeThemeApply:
    def it_copies_colors_and_fonts_from_another_theme(self):
        src_prs = Presentation()
        src_prs.theme.colors[MSO_THEME_COLOR.ACCENT_1] = RGBColor(0x00, 0x80, 0xC0)
        src_prs.theme.fonts.major = "Inter"
        src_prs.theme.fonts.minor = "Inter"

        dst_prs = Presentation()
        dst_prs.theme.apply(src_prs.theme)

        assert dst_prs.theme.colors[MSO_THEME_COLOR.ACCENT_1] == RGBColor(0x00, 0x80, 0xC0)
        assert dst_prs.theme.fonts.major == "Inter"
        assert dst_prs.theme.fonts.minor == "Inter"

    def it_rejects_non_theme_arguments(self, theme):
        with pytest.raises(TypeError):
            theme.apply("not a theme")  # type: ignore[arg-type]


class DescribeThemeName:
    def it_reads_name(self, theme):
        # default theme has a name attribute
        assert isinstance(theme.name, str)

    def it_writes_name(self, theme):
        theme.name = "Brand"
        assert theme.name == "Brand"


class DescribeThemeAccessor:
    """`prs.theme` should be a Theme, not a transient parse."""

    def it_returns_the_same_underlying_element_on_repeated_access(self, prs):
        first = prs.theme
        second = prs.theme
        # Theme objects are recreated on each call (cheap) but they
        # should wrap the *same* underlying element so writes persist.
        assert first._theme_elm is second._theme_elm  # type: ignore[attr-defined]

    def it_exposes_colors_and_fonts(self, theme):
        assert isinstance(theme.colors, ThemeColors)
        assert isinstance(theme.fonts, ThemeFonts)
