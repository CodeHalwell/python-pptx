"""Unit-test suite for high-level `pptx.animation` extensions.

Covers Phase 5 motion paths, the `sequence()` context manager, and
by-paragraph entrance animations.  The lower-level entrance/exit/
emphasis presets are exercised via the integration round-trip suite.
"""

from __future__ import annotations

import pytest

from pptx import Presentation
from pptx.animation import Emphasis, Entrance, MotionPath, Trigger
from pptx.oxml.ns import qn
from pptx.util import Inches


P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


@pytest.fixture
def slide_with_shape():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))
    return slide, shape


def _ctn_node_types(slide):
    return [
        c.get("nodeType")
        for c in slide._element.iter(qn("p:cTn"))
        if c.get("nodeType") is not None
    ]


def _animMotion_paths(slide):
    return [
        m.get("path") for m in slide._element.iter(qn("p:animMotion"))
    ]


class DescribeMotionPath:
    def it_emits_a_path_class_effect_for_a_line(self, slide_with_shape):
        slide, shape = slide_with_shape
        MotionPath.line(slide, shape, Inches(2), Inches(0))
        preset_classes = [
            c.get("presetClass")
            for c in slide._element.iter(qn("p:cTn"))
            if c.get("presetClass") is not None
        ]
        assert preset_classes == ["path"]

    def it_normalizes_the_line_path_against_slide_dimensions(self, slide_with_shape):
        slide, shape = slide_with_shape
        slide_w = slide.part.package.presentation_part.presentation.slide_width
        MotionPath.line(slide, shape, Inches(1), 0)
        path = _animMotion_paths(slide)[0]
        expected_x = float(Inches(1)) / float(slide_w)
        assert path.startswith("M 0 0 L")
        assert path.endswith(" E")
        # Parse the L coordinates: "M 0 0 L X Y E"
        l_part = path.split("L", 1)[1].rsplit("E", 1)[0].strip()
        x_str, y_str = l_part.split()
        assert float(x_str) == pytest.approx(expected_x)
        assert float(y_str) == pytest.approx(0.0)

    def it_emits_a_custom_path_unchanged(self, slide_with_shape):
        slide, shape = slide_with_shape
        MotionPath.custom(slide, shape, "M 0 0 C 0 -0.2 0.2 -0.2 0.2 0 E")
        assert _animMotion_paths(slide) == ["M 0 0 C 0 -0.2 0.2 -0.2 0.2 0 E"]

    def it_rejects_a_malformed_custom_path(self, slide_with_shape):
        slide, shape = slide_with_shape
        with pytest.raises(ValueError):
            MotionPath.custom(slide, shape, "")
        with pytest.raises(ValueError):
            MotionPath.custom(slide, shape, "no end marker")


class DescribeAnimationSequence:
    def it_chains_subsequent_effects_after_the_first(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        a = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))
        b = slide.shapes.add_shape(1, Inches(1), Inches(2.5), Inches(2), Inches(1))
        c = slide.shapes.add_shape(1, Inches(1), Inches(4), Inches(2), Inches(1))

        with slide.animations.sequence():
            Entrance.fade(slide, a)
            Entrance.fade(slide, b)
            Emphasis.pulse(slide, c)

        node_types = [
            nt for nt in _ctn_node_types(slide) if nt != "tmRoot"
        ]
        assert node_types == ["clickEffect", "afterEffect", "afterEffect"]

    def it_honours_an_explicit_start_trigger(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        a = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))
        b = slide.shapes.add_shape(1, Inches(1), Inches(2.5), Inches(2), Inches(1))

        with slide.animations.sequence(start=Trigger.WITH_PREVIOUS):
            Entrance.fade(slide, a)
            Entrance.fade(slide, b)

        node_types = [nt for nt in _ctn_node_types(slide) if nt != "tmRoot"]
        assert node_types == ["withEffect", "afterEffect"]

    def it_does_not_override_an_explicitly_supplied_trigger(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        a = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))
        b = slide.shapes.add_shape(1, Inches(1), Inches(2.5), Inches(2), Inches(1))

        with slide.animations.sequence():
            Entrance.fade(slide, a)
            # Explicit ON_CLICK must win over the sequence's AFTER_PREVIOUS default
            Entrance.fade(slide, b, trigger=Trigger.ON_CLICK)

        node_types = [nt for nt in _ctn_node_types(slide) if nt != "tmRoot"]
        assert node_types == ["clickEffect", "clickEffect"]

    def it_resets_after_the_block(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        a = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(1))
        b = slide.shapes.add_shape(1, Inches(1), Inches(3), Inches(2), Inches(1))

        with slide.animations.sequence():
            Entrance.fade(slide, a)
        # After exit, the next add_* should default back to ON_CLICK
        Entrance.fade(slide, b)

        node_types = [nt for nt in _ctn_node_types(slide) if nt != "tmRoot"]
        assert node_types == ["clickEffect", "clickEffect"]

    def it_rejects_nested_sequences(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        with slide.animations.sequence():
            with pytest.raises(RuntimeError):
                with slide.animations.sequence():
                    pass


class DescribeEntranceByParagraph:
    def it_emits_one_effect_per_paragraph(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        tf = tb.text_frame
        tf.text = "alpha"
        tf.add_paragraph().text = "beta"
        tf.add_paragraph().text = "gamma"

        Entrance.fade(slide, tf, by_paragraph=True)

        # Each paragraph emits a <p:set> + <p:animEffect> pair, both
        # targeting the same paragraph index.  Three paragraphs → three
        # distinct indices, two pRg occurrences per index.
        indices = sorted({
            rg.get("st") for rg in slide._element.iter(qn("p:pRg"))
        })
        assert indices == ["0", "1", "2"]

    def it_chains_paragraphs_after_the_first(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        tf = tb.text_frame
        tf.text = "one"
        tf.add_paragraph().text = "two"

        Entrance.fade(slide, tf, by_paragraph=True)

        node_types = [nt for nt in _ctn_node_types(slide) if nt != "tmRoot"]
        assert node_types == ["clickEffect", "afterEffect"]

    def it_accepts_a_shape_with_a_text_frame(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        tb.text_frame.text = "only"

        Entrance.fade(slide, tb, by_paragraph=True)

        indices = {rg.get("st") for rg in slide._element.iter(qn("p:pRg"))}
        assert indices == {"0"}

    def it_rejects_unsupported_presets(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        tb.text_frame.text = "x"

        with pytest.raises(ValueError):
            slide.animations.add_entrance("fly_in", tb.text_frame, by_paragraph=True)
