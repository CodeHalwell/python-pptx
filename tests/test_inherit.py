"""Unit-test suite for `pptx.inherit`."""

from __future__ import annotations

import pytest

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.inherit import _apply_brightness, resolve_color
from pptx.util import Inches


def _make_run():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    tf = tb.text_frame
    tf.text = "Hello"
    return prs, tf.paragraphs[0].runs[0]


class DescribeResolveColor:
    def it_returns_None_for_an_unset_color(self):
        prs, run = _make_run()
        assert resolve_color(run.font.color, theme=prs.theme) is None

    def it_resolves_an_explicit_RGBColor(self):
        prs, run = _make_run()
        run.font.color.rgb = RGBColor(0x12, 0x34, 0x56)

        assert resolve_color(run.font.color, theme=prs.theme) == RGBColor(0x12, 0x34, 0x56)

    def it_resolves_an_RGBColor_without_a_theme(self):
        prs, run = _make_run()
        run.font.color.rgb = RGBColor(0xFE, 0xDC, 0xBA)

        # `theme` is only required for SCHEME colors.
        assert resolve_color(run.font.color) == RGBColor(0xFE, 0xDC, 0xBA)

    def it_resolves_a_theme_color_via_the_theme(self):
        prs, run = _make_run()
        run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2
        expected = prs.theme.colors[MSO_THEME_COLOR.ACCENT_2]

        assert resolve_color(run.font.color, theme=prs.theme) == expected

    def it_returns_None_for_a_theme_color_when_theme_is_omitted(self):
        prs, run = _make_run()
        run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

        assert resolve_color(run.font.color) is None

    def it_applies_positive_brightness_lightening_toward_white(self):
        prs, run = _make_run()
        run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
        run.font.color.brightness = 0.5  # halfway to white

        # 0x80 + (0xFF - 0x80) * 0.5 = 0x80 + 0x3F.8 ≈ 0xC0
        assert resolve_color(run.font.color, theme=prs.theme) == RGBColor(0xC0, 0xC0, 0xC0)

    def it_applies_negative_brightness_darkening_toward_black(self):
        prs, run = _make_run()
        run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
        run.font.color.brightness = -0.5  # halfway to black

        assert resolve_color(run.font.color, theme=prs.theme) == RGBColor(0x40, 0x40, 0x40)

    def it_returns_None_when_theme_lookup_misses(self):
        """If the theme palette doesn't define a slot the resolver doesn't raise."""
        prs, run = _make_run()
        run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_3

        class _StubTheme:
            class colors:
                @staticmethod
                def __getitem__(key):
                    raise KeyError(key)

        # bind unbound classmethod-style lookup
        stub = _StubTheme()
        stub.colors = type("C", (), {"__getitem__": lambda s, k: (_ for _ in ()).throw(KeyError(k))})()
        assert resolve_color(run.font.color, theme=stub) is None


class DescribeApplyBrightness:
    @pytest.mark.parametrize(
        "brightness,expected",
        [
            (0, RGBColor(0x80, 0x80, 0x80)),
            (None, RGBColor(0x80, 0x80, 0x80)),
            (1.0, RGBColor(0xFF, 0xFF, 0xFF)),
            (-1.0, RGBColor(0x00, 0x00, 0x00)),
            # 0x80 + (0xFF-0x80)*0.25 = 159.75 → rounds to 160 (0xA0)
            (0.25, RGBColor(0xA0, 0xA0, 0xA0)),
            # 0x80 - 0x80*0.25 = 96 (0x60)
            (-0.25, RGBColor(0x60, 0x60, 0x60)),
        ],
    )
    def it_blends_toward_white_or_black(self, brightness, expected):
        assert _apply_brightness(RGBColor(0x80, 0x80, 0x80), brightness) == expected

    def it_clamps_out_of_range_values(self):
        # Values outside [-1, 1] saturate at the endpoints.
        assert _apply_brightness(RGBColor(0x80, 0x80, 0x80), 5.0) == RGBColor(
            0xFF, 0xFF, 0xFF
        )
        assert _apply_brightness(RGBColor(0x80, 0x80, 0x80), -3.0) == RGBColor(
            0x00, 0x00, 0x00
        )
