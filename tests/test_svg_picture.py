"""Integration tests for ``add_svg_picture``.

The cairosvg rasteriser is mocked out in most tests so this suite runs
without the optional dependency installed.
"""

from __future__ import annotations

import io
from unittest.mock import patch

import pytest
from lxml import etree

from pptx import Presentation
from pptx._svg import (
    CairoSvgUnavailable,
    add_svg_blip_extension,
    looks_like_svg,
    rasterize_svg,
)
from pptx.util import Inches


SVG_DOC = b"""<?xml version="1.0" encoding="UTF-8"?>
<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100" viewBox="0 0 100 100">
  <circle cx="50" cy="50" r="40" fill="#FF6600"/>
</svg>"""

# 1x1 transparent PNG
TINY_PNG = bytes.fromhex(
    "89504E470D0A1A0A0000000D49484452000000010000000108060000001F15C4890000"
    "000D49444154789C6364F8FFFF3F0300050001C2DCC1380000000049454E44AE426082"
)


def _new_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    return prs, slide


class DescribeLooksLikeSvg:
    def it_recognises_an_svg_blob(self):
        assert looks_like_svg(SVG_DOC) is True

    def it_rejects_a_png_blob(self):
        assert looks_like_svg(TINY_PNG) is False

    def it_rejects_an_empty_blob(self):
        assert looks_like_svg(b"") is False


class DescribeAddSvgPicture:
    def it_embeds_svg_with_explicit_png_fallback(self):
        prs, slide = _new_slide()

        pic = slide.shapes.add_svg_picture(
            io.BytesIO(SVG_DOC),
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(2),
            png_fallback=TINY_PNG,
        )

        xml = etree.tostring(pic.element).decode()
        assert "svgBlip" in xml
        assert "{96DAC541-7B7A-43D3-8B79-37D633B846F1}" in xml

    def it_round_trips_through_save_and_reopen(self, tmp_path):
        prs, slide = _new_slide()
        slide.shapes.add_svg_picture(
            io.BytesIO(SVG_DOC),
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(2),
            png_fallback=TINY_PNG,
        )

        path = tmp_path / "svg_test.pptx"
        prs.save(str(path))

        reopened = Presentation(str(path))
        # First shape is the title placeholder; the picture is shapes[1].
        pic = reopened.slides[0].shapes[1]
        xml = etree.tostring(pic.element).decode()
        assert "svgBlip" in xml
        assert "{96DAC541-7B7A-43D3-8B79-37D633B846F1}" in xml

    def it_writes_the_svg_part_with_correct_content_type(self, tmp_path):
        prs, slide = _new_slide()
        slide.shapes.add_svg_picture(
            io.BytesIO(SVG_DOC),
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(2),
            png_fallback=TINY_PNG,
        )

        path = tmp_path / "svg_ct.pptx"
        prs.save(str(path))

        import zipfile

        with zipfile.ZipFile(str(path)) as z:
            ct_xml = z.read("[Content_Types].xml").decode()
            names = z.namelist()

        assert "image/svg+xml" in ct_xml
        assert any(n.endswith(".svg") for n in names)

    def it_raises_for_non_svg_input(self):
        prs, slide = _new_slide()

        with pytest.raises(ValueError, match="does not appear to contain SVG"):
            slide.shapes.add_svg_picture(
                io.BytesIO(TINY_PNG),
                Inches(1),
                Inches(1),
                Inches(2),
                Inches(2),
                png_fallback=TINY_PNG,
            )

    def it_uses_cairosvg_when_no_fallback_is_provided(self):
        prs, slide = _new_slide()

        with patch("pptx._svg.rasterize_svg", return_value=TINY_PNG) as raster_:
            slide.shapes.add_svg_picture(
                io.BytesIO(SVG_DOC),
                Inches(1),
                Inches(1),
                Inches(2),
                Inches(2),
            )

        raster_.assert_called_once()


class DescribeRasterizeSvg:
    def it_raises_a_clear_error_when_cairosvg_is_missing(self, monkeypatch):
        # Force the import inside ``rasterize_svg`` to fail.
        import sys

        monkeypatch.setitem(sys.modules, "cairosvg", None)
        with pytest.raises(CairoSvgUnavailable, match="cairosvg"):
            rasterize_svg(SVG_DOC)


class DescribeAddSvgBlipExtension:
    def it_adds_the_extension_to_an_existing_blip(self):
        from pptx.oxml.ns import qn

        # Build a minimal pic element with an <a:blip> child
        pic = etree.fromstring(
            """<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                     xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                     xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
                 <p:blipFill><a:blip r:embed="rId1"/></p:blipFill>
               </p:pic>"""
        )

        add_svg_blip_extension(pic, "rId7")

        ext = pic.find(".//" + qn("a:ext"))
        assert ext is not None
        assert ext.get("uri") == "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
        svgBlip = ext[0]
        assert svgBlip.tag.endswith("}svgBlip")
        assert svgBlip.get(qn("r:embed")) == "rId7"

    def it_raises_when_pic_has_no_blip(self):
        pic = etree.fromstring(
            """<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>"""
        )
        with pytest.raises(ValueError, match="no <a:blip>"):
            add_svg_blip_extension(pic, "rId1")
