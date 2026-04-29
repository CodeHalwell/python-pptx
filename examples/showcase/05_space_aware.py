"""Showcase 05 — Space-aware authoring & the linter.

The whole reason this fork exists.  Three slides:

1. Two boxes side-by-side: the same long string fits one box and
   overflows the other — until ``fit_text`` is called.
2. ``MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE`` as a runtime fallback.
3. A deliberately busted slide that exercises the linter — we
   demonstrate detection then run ``auto_fix`` and re-lint.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from _tokens import BRAND  # noqa: F401  (kept for parity)

HERE = Path(__file__).parent

NEUTRAL = RGBColor(0x0F, 0x17, 0x2A)
PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
SURFACE = RGBColor(0xF8, 0xFA, 0xFC)
DANGER = RGBColor(0xEF, 0x44, 0x44)

LONG_TITLE = (
    "Q4 2026 Customer Outcomes Review — flagship rollouts, "
    "expansion bookings, and headline retention metrics"
)


def build(out_path: Path) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _fit_text_demo(prs)
    _auto_size_demo(prs)
    lint_report = _lint_demo(prs)

    # Final pass: lint the whole deck and auto-fix; the busted slide
    # was deliberately constructed to still raise so we don't call
    # lint_or_die here.  Instead we print the residual report for the
    # reader.
    for slide in prs.slides:
        slide.lint().auto_fix()
    print("\n=== Final lint report (post auto-fix) ===")
    for i, slide in enumerate(prs.slides):
        for issue in slide.lint().issues:
            print(f"slide {i + 1}: [{issue.severity.value}] {issue}")
    print(f"  (initial lint demo found {len(lint_report)} issues "
          f"before auto-fix)")

    prs.save(out_path)
    return prs


def _slide_title(slide, text: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    tf.fit_text(font_family="Inter", max_size=30, bold=True)
    tf.paragraphs[0].font.color.rgb = NEUTRAL

    if subtitle:
        sb = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.15), Inches(12), Inches(0.5),
        )
        sb.text_frame.text = subtitle
        sb.text_frame.paragraphs[0].font.size = Pt(14)
        sb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x64, 0x74, 0x8B)


def _fit_text_demo(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(
        slide,
        "Pre-flight: TextFrame.fit_text",
        "Same long string, two boxes — left is naive, right uses fit_text.",
    )

    # Left: naive Pt(36) — overflows the box.
    # Box geometry deliberately tight (6×1.5") so 36pt overflows
    # massively and fit_text has to drop to ~22pt to make it fit.
    left_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(2.1), Inches(6.0), Inches(1.5),
    )
    _stamp_card(left_box, fill=SURFACE, line=DANGER)
    ltf = left_box.text_frame
    ltf.word_wrap = True
    ltf.margin_left = ltf.margin_right = Pt(14)
    ltf.margin_top = ltf.margin_bottom = Pt(14)
    # ``add_textbox`` defaults to ``SHAPE_TO_FIT_TEXT``, which would
    # silently grow the box and hide the demo. Pin auto_size to NONE
    # so 36pt actually overflows.
    ltf.auto_size = MSO_AUTO_SIZE.NONE
    ltf.text = LONG_TITLE
    # Style the run, not just the paragraph — paragraph-level font
    # properties don't apply to runs that already exist.
    run = ltf.paragraphs[0].runs[0]
    run.font.name = "Inter"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = NEUTRAL
    _caption(slide, Inches(0.6), Inches(5.4), "Naive Pt(36) — overflows the box")

    # Right: identical geometry, fit_text picks the largest whole-pt.
    right_box = slide.shapes.add_textbox(
        Inches(6.7), Inches(2.1), Inches(6.0), Inches(1.5),
    )
    _stamp_card(right_box, fill=SURFACE, line=PRIMARY)
    rtf = right_box.text_frame
    rtf.word_wrap = True
    rtf.margin_left = rtf.margin_right = Pt(14)
    rtf.margin_top = rtf.margin_bottom = Pt(14)
    rtf.text = LONG_TITLE
    # fit_text bakes the largest whole-pt that fits + sets auto_size=NONE.
    rtf.fit_text(font_family="Inter", max_size=36, bold=True)
    rtf.paragraphs[0].runs[0].font.color.rgb = NEUTRAL
    # Belt-and-braces: if the renderer disagrees with fit_text's
    # measurement (it can, by a couple of points, depending on the
    # font fallback), let it shrink further.
    rtf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    _caption(slide, Inches(6.7), Inches(5.4),
             "fit_text — largest whole-pt that fits the box")


def _auto_size_demo(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(
        slide,
        "Runtime: MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE",
        "If the user later types more, PowerPoint shrinks rather than overflows.",
    )

    box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.2), Inches(10.3), Inches(3.4),
    )
    _stamp_card(box, fill=SURFACE, line=PRIMARY)
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = (
        "This text frame has auto_size = TEXT_TO_FIT_SHAPE.\n"
        "Open the deck in PowerPoint, double-click into this box, "
        "and keep typing — the font shrinks instead of overflowing."
    )
    for p in tf.paragraphs:
        p.font.name = "Inter"
        p.font.size = Pt(20)
        p.font.color.rgb = NEUTRAL
    tf.fit_text(font_family="Inter", max_size=22)
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _lint_demo(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(
        slide,
        "Linter: catches what slipped through",
        "An off-slide shape and an overflowing text box — see console output.",
    )

    # Off-slide shape (deliberately past the right edge)
    rogue = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(13.0), Inches(2.2), Inches(2.5), Inches(2.5),
    )
    rogue.fill.solid()
    rogue.fill.fore_color.rgb = DANGER
    rogue.line.fill.background()
    rogue.text_frame.text = "off-slide"
    rogue.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    rogue.text_frame.paragraphs[0].font.bold = True
    rogue.text_frame.paragraphs[0].font.size = Pt(18)

    # Overflowing text frame
    overflow = slide.shapes.add_textbox(
        Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.2),
    )
    _stamp_card(overflow, fill=SURFACE, line=DANGER)
    otf = overflow.text_frame
    otf.word_wrap = True
    otf.text = (
        "This text is set at 28pt without word-wrapping a long enough string "
        "to overflow the height — fit_text would prevent this, but we're "
        "skipping it on purpose to show the lint output."
    )
    otf.paragraphs[0].font.size = Pt(28)
    otf.paragraphs[0].font.color.rgb = NEUTRAL

    # Run the linter and snapshot the issues *before* auto_fix.
    issues = list(slide.lint().issues)
    print("\n=== Lint demo (pre auto-fix) ===")
    for issue in issues:
        print(f"[{issue.severity.value}] {issue}")
    return issues


def _stamp_card(shape, *, fill: RGBColor, line: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)
    shape.shadow.blur_radius = Pt(12)
    shape.shadow.distance = Pt(2)
    shape.shadow.color.rgb = NEUTRAL
    shape.shadow.color.alpha = 0.10


def _caption(slide, left, top, text: str) -> None:
    box = slide.shapes.add_textbox(left, top, Inches(6.0), Inches(0.4))
    box.text_frame.text = text
    p = box.text_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)


if __name__ == "__main__":
    out = HERE / "_out" / "05_space_aware.pptx"
    out.parent.mkdir(exist_ok=True)
    build(out)
    print(f"wrote {out}")
